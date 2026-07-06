#!/usr/bin/env python3
"""批量提取存档文档文本 v1.1 — 支持 .doc/.docx/.pdf/.pptx"""
import os, sys, json, argparse, shutil, subprocess, tempfile
from pathlib import Path

def sanitize_text(text):
    """Remove invalid Unicode surrogates that can appear in PDF extraction."""
    if not isinstance(text, str):
        return ''
    return text.encode('utf-8', errors='replace').decode('utf-8', errors='replace')

def extract_docx(path):
    from docx import Document
    doc = Document(path)
    text = '\n'.join([p.text for p in doc.paragraphs])
    images = sum(1 for rel in doc.part.rels.values() if "image" in rel.reltype)
    domain_codes = sum(1 for p in doc.paragraphs if any(k in p.text for k in ['TOC','HYPERLINK','PAGEREF','REF _Ref','SHAPE']))
    return {'text': text, 'images': images, 'domain_codes': domain_codes}

def extract_doc(path):
    """WPS/Word .doc 二进制提取（三级降级）"""
    # 1. textutil (macOS)
    if shutil.which('textutil'):
        tmp = tempfile.NamedTemporaryFile(suffix='.txt', delete=False)
        tmp.close()
        try:
            r = subprocess.run(['textutil', '-convert', 'txt', '-output', tmp.name, path],
                               capture_output=True, timeout=30)
            if r.returncode == 0:
                with open(tmp.name, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
                if len(text) > 100:
                    return {'text': text, 'method': 'textutil'}
        finally:
            if os.path.exists(tmp.name):
                os.unlink(tmp.name)

    # 2. antiword
    if shutil.which('antiword'):
        r = subprocess.run(['antiword', path], capture_output=True, text=True, timeout=30)
        if r.returncode == 0 and len(r.stdout) > 100:
            return {'text': r.stdout, 'method': 'antiword'}

    # 3. olefile 二进制提取
    try:
        import olefile
        ole = olefile.OleFileIO(path)
        data = ole.openstream('WordDocument').read()
        text_parts = []
        i = 0
        while i < len(data) - 1:
            cp = data[i] | (data[i+1] << 8)
            if (0x4E00 <= cp <= 0x9FFF or 0x20 <= cp <= 0x7E or
                cp in (0x0D, 0x0A, 0xFF0C, 0x3001, 0x3002)):
                start = i
                while i < len(data) - 1:
                    cp2 = data[i] | (data[i+1] << 8)
                    if not (0x4E00 <= cp2 <= 0x9FFF or 0x20 <= cp2 <= 0x7E or
                            cp2 in (0x0D, 0x0A, 0xFF0C, 0x3001, 0x3002)):
                        break
                    i += 2
                chunk = data[start:i].decode('utf-16-le', errors='ignore')
                if len(chunk.strip()) > 3:
                    text_parts.append(chunk)
            i += 1
        ole.close()
        return {'text': '\n'.join(text_parts), 'method': 'olefile'}
    except Exception as e:
        fallback = extract_binary_text(path)
        if fallback:
            return {'text': fallback, 'method': 'binary-fallback', 'warning': str(e)}
        return {'text': '', 'method': 'failed', 'warning': str(e)}

def extract_binary_text(path):
    """Best-effort .doc text fallback when platform converters are unavailable."""
    try:
        data = Path(path).read_bytes()
    except OSError:
        return ''

    chunks = []
    for encoding in ('utf-16-le', 'gb18030', 'utf-8'):
        text = data.decode(encoding, errors='ignore')
        text = ''.join(ch if ch == '\n' or ch == '\r' or ch == '\t' or ord(ch) >= 32 else ' ' for ch in text)
        parts = re_split_text(text)
        chunks.extend(part for part in parts if len(part) >= 8)

    seen = set()
    clean = []
    for chunk in chunks:
        compact = ' '.join(chunk.split())
        if compact and compact not in seen:
            seen.add(compact)
            clean.append(compact)
    return '\n'.join(clean[:200])

def re_split_text(text):
    import re
    return re.findall(r'[\u4e00-\u9fffA-Za-z0-9，。；：、！？,.!?;:\-—（）()《》<>“”"\'\s]{8,}', text)

def extract_pdf(path):
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        from pypdf import PdfReader
    reader = PdfReader(open(path, 'rb'))
    text = ''
    for page in reader.pages:
        t = page.extract_text()
        if t: text += t
    return {'text': sanitize_text(text), 'pages': len(reader.pages)}

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    imgs = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text)
            if shape.shape_type == 13:
                imgs += 1
    return {'text': '\n'.join(texts), 'slides': len(prs.slides), 'images': imgs}

def extract_directory(directory, output):
    os.makedirs(output, exist_ok=True)
    manifest = []

    for root, dirs, files in os.walk(directory):
        for f in sorted(files):
            if f.startswith("~$"): continue
            path = os.path.join(root, f)
            ext = Path(f).suffix.lower()
            try:
                if ext == '.docx': r = extract_docx(path)
                elif ext == '.doc': r = extract_doc(path)
                elif ext == '.pdf': r = extract_pdf(path)
                elif ext == '.pptx': r = extract_pptx(path)
                else: continue

                txt_name = Path(f).stem + '.txt'
                txt_path = os.path.join(output, txt_name)
                with open(txt_path, 'w', encoding='utf-8', errors='replace') as wf:
                    wf.write(sanitize_text(r.get('text', '')))

                entry = {'file': f, 'path': path, 'ext': ext, 'txt_path': txt_path}
                entry.update({k:v for k,v in r.items() if k != 'text'})
                manifest.append(entry)
                method = r.get('method', 'native')
                print(f'OK: {f} ({len(r.get("text",""))} chars, {method})')
            except Exception as e:
                print(f'FAIL: {f} - {e}')

    manifest_path = os.path.join(output, '_manifest.json')
    with open(manifest_path, 'w', encoding='utf-8') as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)
    print(f'\nManifest: {manifest_path} ({len(manifest)} files)')
    return manifest


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('directory', help='存档目录路径')
    parser.add_argument('--output', default='/tmp/auto_grading', help='输出目录')
    args = parser.parse_args()
    extract_directory(args.directory, args.output)

if __name__ == '__main__':
    main()
