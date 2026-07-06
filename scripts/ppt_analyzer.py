#!/usr/bin/env python3
"""
PPT 答辩幻灯片分析器 v1.0 — 全代码模式核心模块

对毕业设计答辩 PPT 进行结构分析和可视化指标提取，
为 PPT 答辩评审提供量化基础数据。

分析维度：
  1. 幻灯片结构分类（封面/目录/背景/设计/成果/总结/致谢）
  2. 图文比例统计
  3. 文字密度评估
  4. 版式一致性检查
  5. 内容完整性检查

输出：PPTReviewResult 包含各维度得分和分析详情

用法:
  python ppt_analyzer.py --pptx demo.pptx --output ppt_analysis.json
"""

import json
import os
import re
import sys
from collections import Counter
from typing import Dict, List, Tuple, Optional, Any


# ============================================================
# PPT 幻灯片加载
# ============================================================
def load_pptx(filepath: str) -> Optional[Any]:
    """加载 PPTX 文件"""
    try:
        from pptx import Presentation
        return Presentation(filepath)
    except ImportError:
        return None
    except Exception as e:
        print(f'[ERROR] 无法加载 {filepath}: {e}')
        return None


# ============================================================
# 幻灯片结构分析器
# ============================================================
class SlideStructureAnalyzer:
    """识别每张幻灯片的类型"""

    # 各类型关键词
    TYPE_KEYWORDS = {
        'cover': ['毕业设计', '毕业论文', '答辩', '题目', '指导教师', '学院',
                  '学号', '姓名', 'graduation', 'thesis', 'defense'],
        'toc': ['目录', 'contents', 'outline', '大纲', '框架', '结构', '流程'],
        'background': ['背景', '现状', '问题', '需求分析', '选题意义', '研究目的',
                       '动机', '意义', 'introduction', 'background', 'motivation'],
        'research': ['文献', '综述', '相关工作', '研究现状', '国内外', '发展',
                     'related work', 'literature', 'survey'],
        'design': ['设计', '架构', '方案', '方法', '框架', '技术', '路线',
                   '模块', '功能', '流程', 'design', 'architecture', 'method'],
        'implementation': ['实现', '开发', '编程', '代码', '算法', '测试',
                           '实验', '数据', '结果', 'implementation', 'code', 'result'],
        'achievement': ['成果', '展示', '演示', '截图', '界面', '效果',
                        'demo', 'showcase', 'achievement', 'output'],
        'conclusion': ['总结', '展望', '结论', '不足', '改进', '收获',
                       '致谢', '感谢', 'conclusion', 'summary', 'future work'],
        'thanks': ['致谢', '感谢', '谢谢', 'Q&A', '问答', '提问',
                   'thanks', 'acknowledgement', 'questions'],
    }

    def __init__(self, config: dict = None):
        self.config = config or {}

    def classify(self, text: str) -> str:
        """
        根据文本内容识别幻灯片类型

        Returns:
            'cover' | 'toc' | 'background' | 'research' | 'design' |
            'implementation' | 'achievement' | 'conclusion' | 'thanks' | 'other'
        """
        text_lower = text.lower()
        scores = {}

        for slide_type, keywords in self.TYPE_KEYWORDS.items():
            score = 0
            for kw in keywords:
                if kw in text_lower:
                    score += 2
            scores[slide_type] = score

        best = max(scores, key=scores.get)
        if scores[best] >= 2:
            return best

        # 位置启发式
        return 'other'

    def analyze_structure(self, slides_data: List[dict]) -> dict:
        """
        分析 PPT 整体结构

        Returns:
            {
                'slide_count': 总页数,
                'type_distribution': 各类型页数分布,
                'structure_score': 结构合理性得分(0-30),
                'missing_types': 缺失的类型,
                'suggestions': 结构建议,
            }
        """
        types = [s['slide_type'] for s in slides_data]
        type_counts = Counter(types)

        # 结构合理性评分
        max_score = 30
        score = max_score
        missing = []
        suggestions = []

        # 必需的幻灯片类型
        required = {
            'cover': '封面页（题目/姓名/导师/学院）',
            'background': '选题背景/研究意义',
            'design': '设计方案/技术路线',
            'achievement': '成果展示/效果演示',
            'conclusion': '总结展望/致谢',
        }

        for req_type, desc in required.items():
            if req_type not in type_counts:
                missing.append(desc)
                score -= 5

        if 'toc' not in type_counts:
            suggestions.append('建议添加目录页，帮助评委了解汇报结构')
            score -= 2

        if 'implementation' not in type_counts and 'design' in type_counts:
            suggestions.append('建议添加实现/测试相关幻灯片，展示技术深度')

        # 页数合理性
        total = len(slides_data)
        if total < 8:
            suggestions.append(f'幻灯片仅{total}页，建议补充至12-20页')
            score -= 3
        elif total > 40:
            suggestions.append(f'幻灯片{total}页较多，建议精简至25页以内')
            score -= 2

        return {
            'slide_count': total,
            'type_distribution': dict(type_counts),
            'structure_score': max(0, score),
            'max_score': max_score,
            'missing_types': missing,
            'suggestions': suggestions,
        }


# ============================================================
# 图文比例分析器
# ============================================================
class ImageTextRatioAnalyzer:
    """分析每页幻灯片的图文比例"""

    def __init__(self, config: dict = None):
        self.config = config or {}

    def analyze(self, slides_data: List[dict]) -> dict:
        """
        分析图文比例

        Returns:
            {
                'avg_text_chars_per_slide': 平均每页字符数,
                'avg_images_per_slide': 平均每页图片数,
                'text_heavy_slides': 文字过多页,
                'image_only_slides': 纯图片页,
                'balance_score': 图文平衡得分(0-25),
            }
        """
        max_score = 25
        score = max_score

        total_chars = sum(s['text_length'] for s in slides_data)
        total_images = sum(s['image_count'] for s in slides_data)
        total_slides = len(slides_data) if slides_data else 1

        avg_chars = total_chars / total_slides
        avg_images = total_images / total_slides

        text_heavy = []
        image_only = []

        for s in slides_data:
            # 每页超过 200 字标记为文字堆砌
            text_len = s['text_length']
            if text_len > 300:
                text_heavy.append({
                    'slide_num': s['slide_num'],
                    'text_length': text_len,
                    'excerpt': s['text'][:50],
                })
                score -= 2
            elif text_len > 200:
                text_heavy.append({
                    'slide_num': s['slide_num'],
                    'text_length': text_len,
                    'excerpt': s['text'][:50],
                })
                score -= 1

            # 纯图片（无文字或极少文字）
            if s['image_count'] >= 1 and text_len < 10:
                image_only.append(s['slide_num'])

        # 理想图文比：有图片的页占40%-80%
        pages_with_images = sum(1 for s in slides_data if s['image_count'] > 0)
        image_ratio = pages_with_images / total_slides if total_slides else 0

        if image_ratio < 0.2:
            score -= 5
        elif image_ratio < 0.4:
            score -= 2

        return {
            'avg_text_chars_per_slide': round(avg_chars, 1),
            'avg_images_per_slide': round(avg_images, 1),
            'text_heavy_slides': text_heavy,
            'image_only_slides': image_only,
            'pages_with_images': pages_with_images,
            'image_ratio': round(image_ratio, 2),
            'balance_score': max(0, score),
            'max_score': max_score,
        }


# ============================================================
# 文字密度评估器
# ============================================================
class TextDensityAnalyzer:
    """评估幻灯片的文字密度和可读性"""

    def __init__(self, config: dict = None):
        self.config = config or {}

    def analyze(self, slides_data: List[dict]) -> dict:
        """
        评估文字密度

        Returns:
            {
                'density_score': 密度合理性得分(0-25),
                'dense_slides': 文字过密页,
                'sparse_slides': 文字过疏页,
            }
        """
        max_score = 25
        score = max_score
        dense = []
        sparse = []

        for s in slides_data:
            text_len = s['text_length']

            if text_len > 350:
                dense.append(f'第{s["slide_num"]}页 ({text_len}字)')
                score -= 3
            elif text_len > 250:
                dense.append(f'第{s["slide_num"]}页 ({text_len}字，偏多)')
                score -= 1

            if text_len < 5 and s['image_count'] == 0:
                sparse.append(f'第{s["slide_num"]}页')
                score -= 1

        return {
            'density_score': max(0, score),
            'max_score': max_score,
            'dense_slides': dense,
            'sparse_slides': sparse,
        }


# ============================================================
# 版式一致性检查器
# ============================================================
class LayoutConsistencyAnalyzer:
    """检查版式一致性"""

    def __init__(self, config: dict = None):
        self.config = config or {}

    def analyze(self, slides_data: List[dict]) -> dict:
        """
        分析版式一致性

        Returns:
            {
                'consistency_score': 版式一致性得分(0-25),
                'issues': 版式问题列表,
            }
        """
        max_score = 25
        score = max_score
        issues = []

        if not slides_data:
            return {'consistency_score': max_score, 'max_score': max_score, 'issues': []}

        # 字体数量统计（如果可获取的话）
        fonts_used = set()
        for s in slides_data:
            for font in s.get('fonts', []):
                fonts_used.add(font)

        if len(fonts_used) > 3:
            score -= min(5, (len(fonts_used) - 3) * 2)
            issues.append(f'使用了 {len(fonts_used)} 种字体（>3种），建议控制在2-3种以内')

        # 幻灯片编号连续性检查
        slide_nums = sorted([s['slide_num'] for s in slides_data])
        if slide_nums[0] != 1:
            issues.append(f'幻灯片编号从{slide_nums[0]}开始而非第1页')
            score -= 1

        # 标题一致性检查
        has_title_count = sum(1 for s in slides_data if s.get('has_title', False))
        title_ratio = has_title_count / len(slides_data) if slides_data else 0

        if title_ratio < 0.5:
            issues.append(f'仅{round(title_ratio*100)}%的幻灯片有标题，建议每页设置标题')
            score -= 3
        elif title_ratio < 0.8:
            issues.append(f'{round(title_ratio*100)}%的幻灯片有标题，少数页面缺标题')
            score -= 1

        return {
            'consistency_score': max(0, score),
            'max_score': max_score,
            'issues': issues,
            'font_count': len(fonts_used),
            'title_coverage': round(title_ratio, 2),
        }


# ============================================================
# PPT 分析器（主入口）
# ============================================================
class PPTAnalyzer:
    """
    PPT 答辩幻灯片分析器 — 全代码模式核心

    对答辩 PPT 进行结构分析和量化评估。
    """

    def __init__(self, config: dict = None):
        self.config = config or {}
        self.structure_analyzer = SlideStructureAnalyzer(config)
        self.ratio_analyzer = ImageTextRatioAnalyzer(config)
        self.density_analyzer = TextDensityAnalyzer(config)
        self.layout_analyzer = LayoutConsistencyAnalyzer(config)

    def analyze(self, pptx_item: dict) -> dict:
        """
        分析 PPT 文件

        Args:
            pptx_item: manifest 中的 PPTX 条目，含 path, file 等

        Returns:
            {
                'ppt_review': {评分+详情},
                'issues': 问题列表,
            }
        """
        pptx_path = pptx_item.get('path', '')
        if not pptx_path or not os.path.exists(pptx_path):
            return {
                'ppt_review': {'error': 'PPT文件不存在'},
                'issues': [{'issue': 'PPT文件不存在', 'confidence': '高'}],
            }

        prs = load_pptx(pptx_path)
        if prs is None:
            return {
                'ppt_review': {'error': 'PPT加载失败'},
                'issues': [{'issue': 'PPT文件加载失败', 'confidence': '高'}],
            }

        # 提取每页幻灯片数据
        slides_data = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            image_count = 0
            fonts = set()
            has_title = False

            for shape in slide.shapes:
                if shape.has_text_frame:
                    frame_text = shape.text_frame.text
                    texts.append(frame_text)

                    # 检测标题
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size and run.font.size >= 200000:  # ≥20pt
                                    has_title = True
                                if run.font.name:
                                    fonts.add(run.font.name)

                if shape.shape_type == 13:  # Picture
                    image_count += 1

                # 检测表格
                if shape.has_table:
                    image_count += 1  # 表格也算作非文本元素

            full_text = ' '.join(texts)
            slides_data.append({
                'slide_num': i,
                'text': full_text,
                'text_length': len(full_text.replace(' ', '').replace('\n', '')),
                'image_count': image_count,
                'fonts': list(fonts),
                'has_title': has_title,
                'slide_type': self.structure_analyzer.classify(full_text),
            })

        # 各维度分析
        structure_result = self.structure_analyzer.analyze_structure(slides_data)
        ratio_result = self.ratio_analyzer.analyze(slides_data)
        density_result = self.density_analyzer.analyze(slides_data)
        layout_result = self.layout_analyzer.analyze(slides_data)

        # 汇总得分（缩放到100分制）
        # 基本信息: 20 | 内容结构: 30 | 视觉美感: 25 | 观点表达: 25
        basic_score = self._calc_basic_score(slides_data, pptx_item)
        structure_score = self._scale(structure_result['structure_score'], structure_result['max_score'], 30)
        visual_score = self._scale(
            ratio_result['balance_score'] + layout_result['consistency_score'],
            ratio_result['max_score'] + layout_result['max_score'],
            25,
        )
        expression_score = self._calc_expression_score(slides_data, density_result)

        ppt_review = {
            '基本信息': basic_score,
            '内容结构': structure_score,
            '视觉美感': visual_score,
            '观点表达': expression_score,
            '总分': basic_score + structure_score + visual_score + expression_score,
        }

        # 汇总问题
        issues = []

        for missing in structure_result.get('missing_types', []):
            issues.append({
                'file': pptx_item.get('file', ''),
                'issue': f'PPT结构缺失: {missing}',
                'confidence': '高',
                'dimension': '内容结构',
            })

        for s in ratio_result.get('text_heavy_slides', []):
            issues.append({
                'file': pptx_item.get('file', ''),
                'issue': f'第{s["slide_num"]}页文字过多({s["text_length"]}字): "{s["excerpt"]}..."',
                'confidence': '高',
                'dimension': '视觉美感',
            })

        for issue in layout_result.get('issues', []):
            issues.append({
                'file': pptx_item.get('file', ''),
                'issue': issue,
                'confidence': '中',
                'dimension': '视觉美感',
            })

        return {
            'ppt_review': ppt_review,
            'structure_analysis': structure_result,
            'ratio_analysis': ratio_result,
            'density_analysis': density_result,
            'layout_analysis': layout_result,
            'issues': issues,
        }

    def _calc_basic_score(self, slides_data: list, pptx_item: dict) -> int:
        """基本信息完整性评分（满分20）"""
        score = 20

        if not slides_data:
            return 0

        # 封面页检测
        cover_text = slides_data[0]['text'] if slides_data else ''

        checks = {
            '题目': r'(题目|标题|课题|毕业设计|论文)',
            '姓名': r'(姓名|学生|作者)',
            '学号': r'\d{6,12}',
            '指导教师': r'(指导教师|导师|advisor)',
            '学院/专业': r'(学院|专业|院系|college|school|department)',
        }

        for field, pattern in checks.items():
            if not re.search(pattern, cover_text, re.IGNORECASE):
                score -= 4

        return max(0, score)

    def _calc_expression_score(self, slides_data: list, density_result: dict) -> int:
        """观点表达评分（满分25）"""
        score = 22  # 基准分

        # 创新点检测
        all_text = ' '.join(s['text'] for s in slides_data)
        innovation_kws = ['创新', '亮点', '特色', '优势', '改进', '优化', '首次', '新颖']
        inno_count = sum(all_text.count(kw) for kw in innovation_kws)

        if inno_count >= 3:
            score += 3
        elif inno_count >= 1:
            score += 1
        else:
            score -= 3

        # 成果展示
        has_demo = any('展示' in s['text'] or '演示' in s['text'] or '效果' in s['text']
                       for s in slides_data)
        has_screenshot = sum(1 for s in slides_data if s['image_count'] > 0) >= 3

        if has_demo and has_screenshot:
            score += 2
        elif not has_demo:
            score -= 2

        # 密度惩罚
        dense_count = len(density_result.get('dense_slides', []))
        score -= min(5, dense_count)

        return max(0, min(25, score))

    def _scale(self, value: float, from_max: float, to_max: float) -> int:
        """分数缩放"""
        if from_max == 0:
            return 0
        return int(value * to_max / from_max)


# ============================================================
# 独立运行支持
# ============================================================
def main():
    import argparse

    parser = argparse.ArgumentParser(description='PPT 答辩幻灯片分析器')
    parser.add_argument('--pptx', required=True, help='PPTX 文件路径')
    parser.add_argument('--output', '-o', help='输出JSON路径')
    args = parser.parse_args()

    analyzer = PPTAnalyzer()

    item = {'path': args.pptx, 'file': os.path.basename(args.pptx)}
    result = analyzer.analyze(item)

    print(json.dumps(result, ensure_ascii=False, indent=2, default=str))

    if args.output:
        with open(args.output, 'w', encoding='utf-8') as f:
            json.dump(result, f, ensure_ascii=False, indent=2, default=str)
        print(f'\n分析结果已保存至: {args.output}')

    ppt_review = result.get('ppt_review', {})
    if ppt_review:
        print(f'\nPPT 评审得分:')
        for dim, s in ppt_review.items():
            if dim != '总分':
                print(f'  {dim}: {s}')
        print(f'  总分: {ppt_review.get("总分", "N/A")}/100')


if __name__ == '__main__':
    main()
